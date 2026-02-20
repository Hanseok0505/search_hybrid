from __future__ import annotations

import csv
import io
import json
import re
import statistics
import uuid
import zipfile
import zlib
from datetime import datetime, timezone
from html import unescape
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from urllib.parse import unquote

import httpx
from fastapi import UploadFile

from app.clients.elasticsearch_client import ElasticsearchClient
from app.clients.graph_client import GraphSearchClient
from app.clients.milvus_client import MilvusVectorClient
from app.core.config import settings
from app.domain.construction import TOP_DOWN_FILTER_KEYS
from app.models.schemas import LLMInvokeRequest
from app.services.embedding_service import EmbeddingService
from app.services.llm_gateway import LLMGateway


class UploadDocumentService:
    def __init__(
        self,
        elastic: Optional[ElasticsearchClient] = None,
        vector: Optional[MilvusVectorClient] = None,
        graph: Optional[GraphSearchClient] = None,
        embedding: Optional[EmbeddingService] = None,
        llm_gateway: Optional[LLMGateway] = None,
    ) -> None:
        self._uploads_dir = Path(settings.uploads_dir)
        self._index_path = Path(settings.uploads_data_path)
        self._uploads_dir.mkdir(parents=True, exist_ok=True)
        self._index_path.parent.mkdir(parents=True, exist_ok=True)
        if not self._index_path.exists():
            self._index_path.write_text("[]", encoding="utf-8")

        self._elastic = elastic
        self._vector = vector
        self._graph = graph
        self._embedding = embedding
        self._llm_gateway = llm_gateway
        self._tika_url = (settings.tika_url or "").strip()
        self._tika_enabled = bool(self._tika_url)
        self._dpe_enabled = bool(settings.upstage_dpe_enabled)
        self._dpe_api_key = (settings.upstage_dpe_api_key or "").strip()
        self._dpe_base_url = (settings.upstage_dpe_base_url or "").strip().rstrip("/")
        self._dpe_timeout_sec = max(5.0, float(settings.upstage_dpe_timeout_sec))
        self._dpe_min_chars = max(0, int(settings.upstage_dpe_min_chars))
        self._auto_context_extract_enabled = bool(settings.auto_context_extract_enabled)
        self._auto_context_extract_use_llm = bool(settings.auto_context_extract_use_llm)
        self._auto_context_extract_max_chars = max(500, int(settings.auto_context_extract_max_chars))
        self._auto_context_extract_min_chars_for_llm = max(
            20,
            int(settings.auto_context_extract_min_chars_for_llm),
        )

    async def ingest(
        self,
        file: UploadFile,
        top_down_context: Optional[Dict[str, str]] = None,
        already_embedded: bool = False,
    ) -> Dict:
        content = await file.read()
        original_file_name = file.filename or "upload.txt"
        content_type = (getattr(file, "content_type", None) or "").strip().lower()
        normalized_name = self._normalize_filename(original_file_name)
        safe_name = self._sanitize_filename(normalized_name)
        safe_name = self._repair_filename_if_needed(safe_name, content_type)
        file_id = str(uuid.uuid4())[:12]
        save_name = f"{file_id}_{safe_name}"
        save_path = self._uploads_dir / save_name
        save_path.write_bytes(content)

        parser = self._detect_parser(safe_name, content_type=content_type)
        warnings_extract: List[str] = []
        text = ""

        if self._prefer_tika_first(parser):
            if self._tika_enabled:
                tika_text = await self._extract_with_tika(content, safe_name, content_type=content_type)
                if tika_text.strip():
                    text = tika_text
                    parser = "tika"
            else:
                warnings_extract.append("tika skipped: TIKA_URL not configured")
        else:
            text = self._extract_text(content, safe_name, content_type=content_type)

        if (not text.strip()) and parser != "tika" and self._tika_enabled:
            tika_text = await self._extract_with_tika(content, safe_name, content_type=content_type)
            if tika_text.strip():
                text = tika_text
                parser = "tika"

        if self._should_try_dpe(parser, safe_name, content_type, text):
            dpe_text, dpe_reason = await self._extract_with_upstage_dpe(
                content=content,
                file_name=safe_name,
                content_type=content_type,
            )
            if dpe_text.strip():
                text = dpe_text
                parser = "upstage_dpe" if parser != "tika" else "tika+upstage_dpe"
            elif dpe_reason:
                warnings_extract.append(f"upstage_dpe skipped: {dpe_reason}")

        if not text.strip():
            native_text = self._extract_text(content, safe_name, content_type=content_type)
            if native_text.strip():
                text = native_text
                if parser == "tika+upstage_dpe":
                    parser = f"{parser}+native"
                elif parser == "tika":
                    parser = "tika+native"
                elif parser == "upstage_dpe":
                    parser = "upstage_dpe+native"
        extracted = bool(text.strip())

        warning: Optional[str] = None
        if not extracted:
            text = f"Uploaded file: {safe_name}"
            warning = "Failed to extract text from file. Please verify parser support for this file type."
        if warnings_extract:
            warning = self._append_warning(warning, " / ".join(warnings_extract))

        manual_context = self._normalize_context(top_down_context or {})
        inferred_context: Dict[str, str] = {}
        llm_context: Dict[str, str] = {}
        extract_sources: list[str] = []

        if self._auto_context_extract_enabled:
            inferred_context = self._infer_context_by_rules(text=text, file_name=safe_name)
            if inferred_context:
                extract_sources.append("rule")
            if self._auto_context_extract_use_llm:
                llm_context, llm_reason = await self._infer_context_with_llm(text=text, file_name=safe_name)
                if llm_context:
                    extract_sources.append("llm")
                elif llm_reason:
                    warning = self._append_warning(warning, f"context llm skipped: {llm_reason}")

        merged_context = self._merge_contexts(
            inferred_context=inferred_context,
            llm_context=llm_context,
            manual_context=manual_context,
        )
        automatic_context_extracted = bool(
            any(k not in manual_context for k in merged_context.keys())
        )

        doc = {
            "id": f"upl-{file_id}",
            "title": safe_name,
            "content": text[:20000],
            "summary": text[:2000],
            "scope_text": "",
            "method_statement": "",
            "risk_register": "",
            "quality_checklist": "",
            "schedule_notes": "",
            "asset_tags": [],
            "wbs_code": "",
            "package_code": "",
            "task_code": "",
            "csi_division": "",
            "spec_section": "",
            "metadata": {
                **merged_context,
                "source": "upload",
                "original_file_name": safe_name,
                "stored_file_name": save_name,
                "uploaded_at": datetime.now(timezone.utc).isoformat(),
                "execution_readiness": 1.0,
                "embedded": bool(already_embedded),
                "parser": parser,
                "context_extract_sources": extract_sources,
                "context_auto_extracted": automatic_context_extracted,
            },
            "related_ids": [],
        }

        index_warnings: List[str] = []
        indexed_backends: List[str] = ["local"]
        embedding_indexed = False

        if self._elastic is not None:
            es_ok, es_reason = await self._safe_index(
                lambda: self._elastic.index_document(doc),
            )
            if es_ok:
                indexed_backends.append("elasticsearch")
            else:
                index_warnings.append(f"elasticsearch indexing skipped: {es_reason}")

        if self._vector is not None:
            vec_text = doc.get("content") or doc.get("summary") or doc["title"]
            if not vec_text.strip():
                vec_text = f"Uploaded file: {safe_name}"
            vec_ok, vec_reason = await self._index_vector(doc, vec_text)
            if vec_ok:
                indexed_backends.append("milvus")
                embedding_indexed = True
                doc["metadata"]["embedded"] = True
            else:
                index_warnings.append(f"milvus indexing skipped: {vec_reason}")

        if self._graph is not None:
            graph_ok, graph_reason = await self._safe_index(
                lambda: self._graph.index_document(doc),
            )
            if graph_ok:
                indexed_backends.append("graph")
            else:
                index_warnings.append(f"graph indexing skipped: {graph_reason}")

        docs = self._load_docs()
        docs.append(doc)
        self._index_path.write_text(json.dumps(docs, ensure_ascii=False, indent=2), encoding="utf-8")

        if index_warnings:
            warning = self._append_warning(warning, " / ".join(index_warnings))

        return {
            "file_name": safe_name,
            "doc_id": doc["id"],
            "indexed": True,
            "chars_indexed": len(doc["content"]),
            "parser": parser,
            "indexed_backends": sorted(set(indexed_backends)),
            "embedding_indexed": embedding_indexed,
            "extracted": extracted,
            "automatic_context_extracted": automatic_context_extracted,
            "top_down_context": merged_context,
            "warning": warning,
        }

    @staticmethod
    def _prefer_tika_first(parser: str) -> bool:
        # Keep lightweight text-like formats on native parser to avoid network round-trip.
        return parser not in {"text", "html"}

    def _should_try_dpe(
        self,
        parser: str,
        file_name: str,
        content_type: str,
        extracted_text: str,
    ) -> bool:
        suffix = Path(file_name).suffix.lower().lstrip(".")
        advanced_ext = {
            "pdf",
            "doc",
            "docx",
            "xls",
            "xlsx",
            "ppt",
            "pptx",
            "hwp",
            "hwpx",
            "png",
            "jpg",
            "jpeg",
            "tif",
            "tiff",
        }
        looks_advanced = (
            suffix in advanced_ext
            or parser in {"pdf", "doc", "docx", "excel", "pptx", "hwp", "hwpx", "fallback", "tika"}
            or "pdf" in content_type
            or "officedocument" in content_type
            or "msword" in content_type
            or "powerpoint" in content_type
            or "excel" in content_type
        )
        if not looks_advanced:
            return False
        return len((extracted_text or "").strip()) < self._dpe_min_chars

    def _is_dpe_available(self) -> tuple[bool, str]:
        if not self._dpe_enabled:
            return False, "UPSTAGE_DPE_ENABLED=false"
        if not self._dpe_api_key:
            return False, "UPSTAGE_DPE_API_KEY missing"
        if not self._dpe_base_url:
            return False, "UPSTAGE_DPE_BASE_URL missing"
        return True, ""

    async def _extract_with_upstage_dpe(
        self,
        *,
        content: bytes,
        file_name: str,
        content_type: str = "",
    ) -> tuple[str, str]:
        available, reason = self._is_dpe_available()
        if not available:
            return "", reason
        if not content:
            return "", "empty content"

        url = self._dpe_base_url
        headers = {
            "Authorization": f"Bearer {self._dpe_api_key}",
            "Accept": "application/json",
        }
        payload_type = content_type or "application/octet-stream"
        attempts: list[tuple[str, dict]] = [
            ("document", {"document": (file_name, content, payload_type)}),
            ("file", {"file": (file_name, content, payload_type)}),
        ]

        async with httpx.AsyncClient(timeout=self._dpe_timeout_sec) as client:
            last_reason = "unknown dpe error"
            for field_name, files in attempts:
                try:
                    response = await client.post(url, headers=headers, files=files)
                except Exception as exc:
                    last_reason = str(exc)
                    continue

                if response.status_code >= 400:
                    body = (response.text or "").strip().replace("\n", " ")
                    last_reason = f"status {response.status_code}: {body[:220]}"
                    continue

                try:
                    payload = response.json()
                except Exception:
                    payload = {"text": response.text}

                text = self._extract_text_from_dpe_payload(payload)
                if text.strip():
                    return text.strip(), ""
                last_reason = f"empty response on field '{field_name}'"

            return "", last_reason

    @staticmethod
    def _extract_text_from_dpe_payload(payload: object) -> str:
        chunks: list[str] = []

        def _push(value: object) -> None:
            if isinstance(value, str):
                stripped = value.strip()
                if stripped:
                    chunks.append(stripped)

        if isinstance(payload, str):
            _push(payload)
            return "\n".join(chunks)

        if not isinstance(payload, dict):
            return ""

        for key in ("content", "text", "markdown", "output"):
            _push(payload.get(key))

        result = payload.get("result")
        if isinstance(result, dict):
            for key in ("content", "text", "markdown", "output"):
                _push(result.get(key))

        pages = payload.get("pages")
        if isinstance(pages, list):
            for page in pages:
                if isinstance(page, dict):
                    for key in ("content", "text", "markdown", "output"):
                        _push(page.get(key))

        elements = payload.get("elements")
        if isinstance(elements, list):
            for item in elements:
                if isinstance(item, dict):
                    for key in ("content", "text", "markdown", "output"):
                        _push(item.get(key))

        return "\n".join(chunks)

    @staticmethod
    def _normalize_context(raw_context: Dict[str, object]) -> Dict[str, str]:
        out: Dict[str, str] = {}
        for key, value in (raw_context or {}).items():
            if key not in TOP_DOWN_FILTER_KEYS:
                continue
            if value is None:
                continue
            text = str(value).strip()
            if not text:
                continue
            out[key] = text
        return out

    def _merge_contexts(
        self,
        *,
        inferred_context: Dict[str, str],
        llm_context: Dict[str, str],
        manual_context: Dict[str, str],
    ) -> Dict[str, str]:
        merged: Dict[str, str] = {}
        merged.update(self._normalize_context(inferred_context))
        merged.update(self._normalize_context(llm_context))
        merged.update(self._normalize_context(manual_context))
        return merged

    def _infer_context_by_rules(self, *, text: str, file_name: str) -> Dict[str, str]:
        source = f"{file_name}\n{text[: self._auto_context_extract_max_chars]}"
        normalized_upper = source.upper()
        normalized_search = re.sub(r"[_/]+", " ", normalized_upper)
        out: Dict[str, str] = {}

        project_id = self._regex_group(
            normalized_search,
            r"\b(PROJ-[A-Z0-9]+(?:-[A-Z0-9]+){0,4})\b",
            1,
        )
        if project_id:
            out["project_id"] = project_id.replace("_", "-")

        tower_direct = self._regex_group(
            normalized_search,
            r"\b(TOWER[-_ ]?[A-Z0-9]{1,6})\b",
            1,
        )
        if tower_direct:
            out["building"] = re.sub(r"[-_ ]+", "-", tower_direct)
        else:
            building_code = self._regex_group(
                normalized_search,
                r"\b(?:BUILDING|BLDG)[-_ ]?(?:TOWER[-_ ]?)?([A-Z0-9]{1,6})\b",
                1,
            )
            if building_code:
                out["building"] = f"TOWER-{building_code}"

        level = self._regex_group(
            normalized_search,
            r"\b(B\d{1,2}|L\d{1,2}|F\d{1,2}|RF|ROOF|\d{1,2}F)\b",
            1,
        )
        if level:
            out["level"] = level

        package_code = self._regex_group(normalized_search, r"\b(PKG[-_][A-Z0-9-]{2,})\b", 1)
        if package_code:
            out["package_code"] = package_code.replace("_", "-")

        task_code = self._regex_group(normalized_search, r"\b(TASK[-_][A-Z0-9-]{2,})\b", 1)
        if task_code:
            out["task_code"] = task_code.replace("_", "-")

        wbs_suffix = self._regex_group(normalized_search, r"\bWBS[-_:. ]?([A-Z0-9.-]{2,})\b", 1)
        if wbs_suffix:
            wbs_suffix = re.sub(r"^WBS[-_:. ]*", "", wbs_suffix, flags=re.IGNORECASE)
        wbs_code = f"WBS-{wbs_suffix}" if wbs_suffix else ""
        if wbs_code:
            out["wbs_code"] = (
                wbs_code.replace("_", "-")
                .replace(" ", "")
                .replace(":", "-")
                .replace("--", "-")
            )

        csi_match = re.search(r"\b(?:CSI|DIVISION)[-_ ]?(\d{2})\b", normalized_search)
        if csi_match:
            out["csi_division"] = csi_match.group(1)

        spec_section = self._regex_group(
            normalized_search,
            r"\b(?:SPEC|SECTION)[-_ ]?(\d{2}[-_ ]?\d{2}[-_ ]?\d{2}|\d{6})\b",
            1,
        )
        if spec_section:
            out["spec_section"] = spec_section.replace("_", "-").replace(" ", "-")

        work_type_map = {
            "concrete": ["CONCRETE", "SLAB", "POUR", "FORMWORK", "REBAR"],
            "earthwork": ["EXCAVATION", "EARTHWORK", "DEWATERING"],
            "mep": ["MEP", "HVAC", "DUCT", "PIPING", "PLUMBING", "ELECTRICAL"],
            "facade": ["FACADE", "CURTAIN WALL", "CLADDING", "ENVELOPE"],
            "safety": ["SAFETY", "HSE", "EHS", "PERMIT TO WORK"],
        }
        for work_type, keywords in work_type_map.items():
            if any(keyword in normalized_upper for keyword in keywords):
                out["work_type"] = work_type
                break

        return self._normalize_context(out)

    async def _infer_context_with_llm(self, *, text: str, file_name: str) -> tuple[Dict[str, str], str]:
        if not self._llm_gateway:
            return {}, "llm gateway unavailable"

        content = (text or "").strip()
        if len(content) < self._auto_context_extract_min_chars_for_llm:
            return {}, "text too short for llm context extraction"

        prompt = (
            "Extract construction top-down metadata from the provided filename and text.\n"
            "Return STRICT JSON object only. No markdown.\n"
            f"Allowed keys: {', '.join(TOP_DOWN_FILTER_KEYS)}\n"
            "Use only explicit evidence. If unknown, omit the key.\n\n"
            f"filename: {file_name}\n"
            f"text:\n{content[: self._auto_context_extract_max_chars]}"
        )

        try:
            out = await self._llm_gateway.invoke(
                LLMInvokeRequest(
                    task="chat",
                    provider=settings.llm_provider,
                    model=settings.auto_context_extract_llm_model or settings.llm_rerank_model,
                    text=prompt,
                    temperature=0.0,
                    max_tokens=420,
                )
            )
            raw_text = str(out.get("output_text") or "").strip()
            payload = self._parse_json_object(raw_text)
            if not isinstance(payload, dict):
                return {}, "invalid llm json payload"
            normalized = self._normalize_context(payload)
            if not normalized:
                return {}, "llm returned empty context"
            return normalized, ""
        except Exception as exc:
            return {}, str(exc)

    @staticmethod
    def _parse_json_object(raw_text: str) -> Dict[str, object]:
        text = (raw_text or "").strip()
        if not text:
            return {}

        try:
            parsed = json.loads(text)
            return parsed if isinstance(parsed, dict) else {}
        except Exception:
            pass

        start = text.find("{")
        end = text.rfind("}")
        if start == -1 or end <= start:
            return {}
        snippet = text[start : end + 1]
        try:
            parsed = json.loads(snippet)
            return parsed if isinstance(parsed, dict) else {}
        except Exception:
            return {}

    @staticmethod
    def _regex_group(source_text: str, pattern: str, group: int) -> str:
        match = re.search(pattern, source_text, flags=re.IGNORECASE)
        if not match:
            return ""
        value = (match.group(group) or "").strip()
        return value.upper()

    def _append_warning(self, base: Optional[str], extra: str) -> str:
        if not base:
            return extra
        return f"{base} / {extra}"

    @staticmethod
    def _normalize_filename(file_name: str) -> str:
        file_name = file_name.strip()
        if not file_name:
            return "upload.txt"
        if file_name.lower().startswith("utf-8''"):
            file_name = unquote(file_name[7:], errors="replace")
        try:
            decoded = unquote(file_name)
        except Exception:
            decoded = file_name
        decoded = UploadDocumentService._repair_filename_encoding(decoded)
        return decoded

    @staticmethod
    def _repair_filename_encoding(file_name: str) -> str:
        if not file_name:
            return file_name
        if not any(ord(ch) > 127 for ch in file_name):
            return file_name

        # Handle mojibake like UTF-8 bytes decoded as cp1252/latin1.
        try:
            latin1_bytes = file_name.encode("latin-1", errors="strict")
        except Exception:
            return file_name

        for encoding in ("utf-8", "cp949", "utf-16le", "utf-16be"):
            try:
                repaired = latin1_bytes.decode(encoding, errors="strict")
                if repaired != file_name:
                    return repaired
            except Exception:
                continue
        return file_name

    @staticmethod
    def _sanitize_filename(file_name: str) -> str:
        base = re.sub(r"[\\/:*?\"<>|]", "_", file_name, flags=re.UNICODE)
        base = re.sub(r"\s+", "_", base, flags=re.UNICODE).strip("._ ")
        base = re.sub(r"_{2,}", "_", base)
        return base or "upload.txt"

    @staticmethod
    def _infer_ext_from_content_type(content_type: str) -> str:
        ct = (content_type or "").lower()
        if not ct:
            return ""
        if "plain" in ct or ct == "text/plain":
            return "txt"
        if "json" in ct:
            return "json"
        if "csv" in ct:
            return "csv"
        if "pdf" in ct:
            return "pdf"
        if "wordprocessingml.document" in ct:
            return "docx"
        if "msword" in ct:
            return "doc"
        if "presentationml.presentation" in ct or "powerpoint" in ct:
            return "pptx"
        if "spreadsheetml.sheet" in ct or "excel" in ct:
            return "xlsx"
        if "html" in ct:
            return "html"
        if "hwp" in ct:
            return "hwp"
        return ""

    @staticmethod
    def _repair_filename_if_needed(file_name: str, content_type: str) -> str:
        candidate = file_name.strip()
        lower = candidate.lower()
        stem = Path(candidate).stem
        suffix = Path(candidate).suffix.lower().lstrip(".")

        known_exts = {
            "txt", "md", "csv", "json", "log", "yaml", "yml", "pdf",
            "docx", "doc", "xlsx", "xls", "pptx", "html", "htm", "hwpx", "hwp",
        }
        if (
            lower in known_exts
            or (not suffix and len(stem) <= 1)
            or suffix not in known_exts
            or candidate == "upload.txt"
        ):
            inferred = UploadDocumentService._infer_ext_from_content_type(content_type)
            if inferred:
                if lower.endswith(f".{inferred}"):
                    return candidate
                if candidate.lower() in known_exts:
                    return f"upload.{inferred}"
                if not suffix:
                    return f"{candidate}.{inferred}"
                return f"upload.{inferred}"

        return candidate

    def _load_docs(self) -> List[Dict]:
        try:
            return json.loads(self._index_path.read_text(encoding="utf-8"))
        except Exception:
            return []

    @staticmethod
    def _detect_parser(file_name: str, content_type: str = "") -> str:
        lower = file_name.lower()
        if lower.endswith((".txt", ".md", ".csv", ".json", ".log", ".yaml", ".yml")):
            return "text"
        if lower.endswith(".pdf"):
            return "pdf"
        if lower.endswith(".docx"):
            return "docx"
        if lower.endswith(".doc"):
            return "doc"
        if lower.endswith(".xlsx") or lower.endswith(".xls"):
            return "excel"
        if lower.endswith(".pptx"):
            return "pptx"
        if lower.endswith((".html", ".htm")):
            return "html"
        if lower.endswith(".hwpx"):
            return "hwpx"
        if lower.endswith(".hwp"):
            return "hwp"
        if lower.startswith("."):
            inferred = UploadDocumentService._infer_ext_from_content_type(content_type)
            if inferred in {"txt", "md", "csv", "json", "pdf", "doc", "docx", "xlsx", "xls", "pptx", "html", "hwp", "hwpx"}:
                if inferred == "txt":
                    return "text"
                if inferred in {"xls", "xlsx"}:
                    return "excel"
                return inferred
        return "fallback"

    @staticmethod
    def _detect_encoding_candidates() -> Tuple[str, ...]:
        return ("utf-8", "cp949", "euc-kr", "utf-16", "utf-16le", "utf-16be")

    @staticmethod
    def _decode_bytes(content: bytes) -> str:
        # UTF-16 without BOM can incorrectly decode CP949 bytes.
        # Keep UTF-8/legacy Korean encodings first and then UTF-16 as fallback.
        for enc in UploadDocumentService._detect_encoding_candidates():
            try:
                text = content.decode(enc)
                if UploadDocumentService._is_reasonable_text(text):
                    return text
            except Exception:
                continue
        fallback = content.decode("utf-8", errors="ignore")
        if UploadDocumentService._is_reasonable_text(fallback):
            return fallback
        return ""

    @staticmethod
    def _is_reasonable_text(text: str) -> bool:
        if not text:
            return False
        stripped = text.replace("\x00", "")
        if len(stripped) < 2:
            return False
        visible = [ch for ch in stripped if ch.isprintable() and ch not in {"\u200b", "\ufeff"}]
        ratio = len(visible) / len(stripped)
        # binary gibberish often has very low printable ratio.
        if ratio < 0.7:
            return False
        # reject almost all control-characters payload.
        control_ratio = statistics.mean(
            1.0 if ch in {"\n", "\r", "\t"} or ord(ch) >= 32 else 0.0 for ch in stripped
        )
        return control_ratio > 0.85

    @staticmethod
    def _extract_text(content: bytes, file_name: str, content_type: str = "") -> str:
        lower = file_name.lower()
        if (not lower or lower.startswith(".")) and content_type:
            inferred = UploadDocumentService._infer_ext_from_content_type(content_type)
            if inferred:
                if inferred == "txt":
                    lower = ".txt"
                elif inferred == "json":
                    lower = ".json"
                elif inferred == "csv":
                    lower = ".csv"
                elif inferred == "pdf":
                    lower = ".pdf"
                elif inferred == "doc":
                    lower = ".doc"
                elif inferred == "docx":
                    lower = ".docx"
                elif inferred in {"xls", "xlsx"}:
                    lower = ".xlsx"
                elif inferred == "pptx":
                    lower = ".pptx"
                elif inferred == "html":
                    lower = ".html"
                elif inferred == "hwp":
                    lower = ".hwp"
        if lower.endswith((".txt", ".md", ".log", ".yaml", ".yml")):
            return UploadDocumentService._decode_bytes(content).strip()
        if lower.endswith((".csv",)):  # csv to row-major text preserve structure.
            return UploadDocumentService._extract_csv_text(content)
        if lower.endswith((".json",)):
            return UploadDocumentService._extract_json_text(content)
        if lower.endswith(".pdf"):
            return UploadDocumentService._extract_pdf_text(content)
        if lower.endswith(".docx"):
            return UploadDocumentService._extract_docx_text(content)
        if lower.endswith(".doc"):
            return UploadDocumentService._extract_doc_legacy_text(content)
        if lower.endswith(".pptx"):
            return UploadDocumentService._extract_pptx_text(content)
        if lower.endswith(".xlsx"):
            return UploadDocumentService._extract_xlsx_text(content)
        if lower.endswith(".xls"):
            return UploadDocumentService._extract_xls_text(content)
        if lower.endswith((".html", ".htm")):
            return UploadDocumentService._extract_html_text(content)
        if lower.endswith(".hwpx"):
            return UploadDocumentService._extract_hwpx_text(content)
        if lower.endswith(".hwp"):
            return UploadDocumentService._extract_hwp_text(content)

        return UploadDocumentService._decode_bytes(content)

    @staticmethod
    def _extract_csv_text(content: bytes) -> str:
        raw = UploadDocumentService._decode_bytes(content)
        try:
            stream = io.StringIO(raw)
            reader = csv.reader(stream)
            rows = []
            for row in reader:
                row_values = [str(item).strip() for item in row if str(item).strip()]
                if row_values:
                    rows.append(" | ".join(row_values))
            return "\n".join(rows)
        except Exception:
            return raw

    @staticmethod
    def _extract_json_text(content: bytes) -> str:
        raw = UploadDocumentService._decode_bytes(content)
        try:
            obj = json.loads(raw)
        except Exception:
            return raw

        parts = []
        if isinstance(obj, dict):
            for key, value in obj.items():
                parts.append(f"{key}: {value}")
        elif isinstance(obj, list):
            for item in obj[:200]:
                parts.append(str(item))
        else:
            parts.append(str(obj))
        return "\n".join(parts)

    @staticmethod
    def _extract_pdf_text(content: bytes) -> str:
        try:
            import fitz  # type: ignore
        except Exception:
            return ""

        try:
            texts: List[str] = []
            with fitz.open(stream=content, filetype="pdf") as doc:
                for page in doc:
                    page_text = page.get_text("text")
                    if page_text:
                        texts.append(page_text)
            return "\n".join(texts)
        except Exception:
            return ""

    @staticmethod
    def _extract_docx_text(content: bytes) -> str:
        try:
            from docx import Document  # type: ignore
        except Exception:
            return ""

        try:
            doc = Document(io.BytesIO(content))
            chunks: List[str] = []
            for para in doc.paragraphs:
                text = (para.text or "").strip()
                if text:
                    chunks.append(text)
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    row_text = " | ".join(value for value in cells if value)
                    if row_text:
                        chunks.append(row_text)
            return "\n".join(chunks)
        except Exception:
            return ""

    @staticmethod
    def _extract_doc_legacy_text(content: bytes) -> str:
        # Best-effort for legacy binary DOC or unknown legacy word formats.
        return UploadDocumentService._decode_bytes(content)

    @staticmethod
    def _extract_pptx_text(content: bytes) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except Exception:
            return ""

        try:
            prs = Presentation(io.BytesIO(content))
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    try:
                        text = (shape.text or "").strip()
                        if text:
                            out.append(text)
                    except Exception:
                        continue
            return "\n".join(out)
        except Exception:
            return ""

    @staticmethod
    def _extract_xlsx_text(content: bytes) -> str:
        try:
            import openpyxl  # type: ignore
        except Exception:
            return ""

        try:
            wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets:
                out.append(f"[sheet]{ws.title}")
                for row in ws.iter_rows(values_only=True):
                    text = " | ".join(str(item).strip() for item in row if item is not None)
                    if text:
                        out.append(text)
            return "\n".join(out)
        except Exception:
            return ""

    @staticmethod
    def _extract_xls_text(content: bytes) -> str:
        # Legacy excel (.xls) requires additional optional libs. Best-effort fallback.
        return UploadDocumentService._decode_bytes(content)

    @staticmethod
    def _extract_html_text(content: bytes) -> str:
        raw = UploadDocumentService._decode_bytes(content)
        raw = re.sub(r"<script[\s\S]*?</script>", " ", raw, flags=re.IGNORECASE)
        raw = re.sub(r"<style[\s\S]*?</style>", " ", raw, flags=re.IGNORECASE)
        stripped = re.sub(r"<[^>]+>", " ", raw)
        return re.sub(r"\s+", " ", unescape(stripped)).strip()

    @staticmethod
    def _extract_hwpx_text(content: bytes) -> str:
        texts: List[str] = []
        try:
            with zipfile.ZipFile(io.BytesIO(content)) as zf:
                names = [n for n in zf.namelist() if n.lower().endswith(".xml")]
                for name in names:
                    lower_name = name.lower()
                    if "contents" not in lower_name and "section" not in lower_name:
                        continue
                    raw = zf.read(name)
                    xml = raw.decode("utf-8", errors="ignore")
                    stripped = re.sub(r"<[^>]+>", " ", xml)
                    clean = re.sub(r"\s+", " ", unescape(stripped)).strip()
                    if clean:
                        texts.append(clean)
        except Exception:
            return ""
        return "\n".join(texts)[:300000]

    @staticmethod
    def _extract_hwp_text(content: bytes) -> str:
        try:
            import olefile  # type: ignore
        except Exception:
            return ""

        try:
            ole = olefile.OleFileIO(io.BytesIO(content))
        except Exception:
            return ""

        try:
            if ole.exists("PrvText"):
                raw = ole.openstream("PrvText").read()
                for enc in ("utf-16le", "cp949", "utf-8"):
                    try:
                        txt = raw.decode(enc, errors="ignore").strip()
                        if txt:
                            return txt[:300000]
                    except Exception:
                        continue

            sections = []
            for path in ole.listdir():
                if len(path) == 2 and path[0] == "BodyText" and path[1].startswith("Section"):
                    sections.append(path)
            sections.sort(key=lambda p: p[1])

            out: List[str] = []
            for sec in sections:
                data = ole.openstream(sec).read()
                try:
                    data = zlib.decompress(data, -15)
                except Exception:
                    pass

                pos = 0
                while pos + 4 <= len(data):
                    header = int.from_bytes(data[pos : pos + 4], "little")
                    pos += 4
                    rec_type = header & 0x3FF
                    rec_len = (header >> 20) & 0xFFF
                    if pos + rec_len > len(data):
                        break
                    payload = data[pos : pos + rec_len]
                    pos += rec_len

                    if rec_type == 67 and payload:
                        txt = payload.decode("utf-16le", errors="ignore")
                        txt = re.sub(r"\s+", " ", txt).strip()
                        if txt:
                            out.append(txt)

            return "\n".join(out)[:300000]
        except Exception:
            return ""
        finally:
            try:
                ole.close()
            except Exception:
                pass

    async def _extract_with_tika(
        self,
        content: bytes,
        file_name: str,
        content_type: str = "",
    ) -> str:
        if not self._tika_enabled or not content:
            return ""
        base_url = self._tika_url.rstrip("/")
        payload_type = content_type or "application/octet-stream"
        headers = {"Accept": "text/plain"}

        attempts: list[tuple[str, str, object, dict]] = [
            (
                "post",
                f"{base_url}/tika/form",
                {"file": (file_name, content, payload_type)},
                {"timeout": 30.0},
            ),
            (
                "put",
                f"{base_url}/tika",
                content,
                {"timeout": 30.0, "headers": {"Content-Type": payload_type}},
            ),
            (
                "post",
                f"{base_url}/rmeta/form",
                {"file": (file_name, content, payload_type)},
                {"timeout": 30.0},
            ),
        ]

        async with httpx.AsyncClient(timeout=20.0) as client:
            for method, endpoint, body, options in attempts:
                try:
                    if method == "post":
                        response = await client.post(
                            endpoint,
                            data=body if isinstance(body, dict) and "file" not in body else None,
                            files=body if isinstance(body, dict) else None,
                            headers=headers,
                            timeout=options.get("timeout", 20.0),
                        )
                    else:
                        response = await client.put(
                            endpoint,
                            content=body if isinstance(body, (bytes, bytearray)) else None,
                            headers={**headers, **options.get("headers", {})},
                            timeout=options.get("timeout", 20.0),
                        )
                except Exception:
                    continue

                if response.status_code >= 400:
                    continue

                text = (response.text or "").strip()
                if text:
                    return text

        return ""

    async def _safe_index(self, func) -> Tuple[bool, str]:
        try:
            result = func()
            if hasattr(result, "__await__"):
                ok = await result
            else:
                ok = result
            return bool(ok), ""
        except Exception as exc:  # pragma: no cover - defensive
            return False, str(exc)

    async def _index_vector(self, doc: Dict, content: str) -> Tuple[bool, str]:
        if self._vector is None:
            return False, "backend unavailable"
        if self._embedding is None:
            return False, "embedding client unavailable"

        try:
            embedding = await self._embedding.embed(content[:4000])
        except Exception as exc:
            return False, f"embedding failed: {exc}"

        if hasattr(self._vector, "index_document_with_reason"):
            reason = self._vector.index_document_with_reason(doc, embedding)
            if isinstance(reason, tuple) and len(reason) == 2:
                return bool(reason[0]), str(reason[1] or "")

        return await self._safe_index(lambda: self._vector.index_document(doc, embedding))

